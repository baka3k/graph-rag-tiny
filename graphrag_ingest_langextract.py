import argparse
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Tuple

from neo4j import GraphDatabase
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.http import models as qmodels
from sentence_transformers import SentenceTransformer

from embedding_utils import resolve_embedding_model

try:
    from dotenv import load_dotenv
except Exception:
    load_dotenv = None

from entity_extractors import (
    build_spacy_pipeline,
    extract_entities_gemini,
    extract_entities_gliner,
    extract_entities_gliner_batch,
    extract_entities_langextract,
    build_gliner_model,
    parse_gliner_labels,
)


def _load_env() -> None:
    if load_dotenv is None:
        return
    env_path = Path(__file__).with_name(".env")
    if env_path.exists():
        load_dotenv(dotenv_path=env_path)
    else:
        load_dotenv()


def read_pdf_text(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n\n".join(p for p in parts if p.strip())


def read_text_file(text_path: Path) -> str:
    return text_path.read_text(encoding="utf-8").strip()


def read_docx_text(docx_path: Path) -> str:
    doc = Document(str(docx_path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(parts)


def read_pptx_text(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                stripped = text.strip()
                if stripped:
                    parts.append(stripped)
    return "\n\n".join(parts)


def split_paragraphs(text: str, max_chars: int) -> List[str]:
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n+", text) if p.strip()]
    chunks: List[str] = []
    for para in paragraphs:
        if len(para) <= max_chars:
            chunks.append(para)
        else:
            for i in range(0, len(para), max_chars):
                chunk = para[i : i + max_chars].strip()
                if chunk:
                    chunks.append(chunk)
    return chunks


def _set_langextract_overrides(model_id: str | None, model_url: str | None) -> None:
    if model_id:
        os.environ["LANGEXTRACT_MODEL_ID"] = model_id
    if model_url:
        os.environ["LANGEXTRACT_MODEL_URL"] = model_url


def _normalize_entity_name(name: str) -> str:
    return re.sub(r"\s+", " ", name.strip().lower())


def _entity_id(name: str, ent_type: str) -> str:
    key = f"{ent_type}::{_normalize_entity_name(name)}"
    return str(uuid.uuid5(uuid.NAMESPACE_URL, key))


def build_graph_components(
    text: str,
    provider: str,
    nlp=None,
    gliner_model=None,
    gliner_labels: List[str] | None = None,
    gliner_threshold: float = 0.3,
    merge_entities: bool = True,
) -> Tuple[Dict[str, Dict[str, str]], List[Dict[str, str]]]:
    if provider == "spacy":
        doc = nlp(text)
        entities = [{"name": ent.text, "type": ent.label_} for ent in doc.ents]
        relations = []
    elif provider == "gliner":
        entities, relations = extract_entities_gliner(
            text,
            labels=gliner_labels,
            threshold=gliner_threshold,
            gliner_model=gliner_model,
        )
    elif provider == "gemini":
        entities, relations = extract_entities_gemini(text)
    else:
        entities, relations = extract_entities_langextract(text)
    return build_graph_components_from_entities(entities, relations, merge_entities=merge_entities)


def build_graph_components_from_entities(
    entities: List[Dict[str, str]],
    relations: List[Dict[str, str]],
    merge_entities: bool = True,
) -> Tuple[Dict[str, Dict[str, str]], List[Dict[str, str]]]:
    nodes: Dict[str, Dict[str, str]] = {}
    name_index: Dict[str, str] = {}
    for ent in entities:
        name = ent.get("name", "").strip()
        if not name:
            continue
        ent_type = ent.get("type", "UNKNOWN") or "UNKNOWN"
        name_norm = _normalize_entity_name(name)
        if merge_entities:
            key = f"{ent_type}::{name_norm}"
            if key not in nodes:
                nodes[key] = {
                    "id": _entity_id(name, ent_type),
                    "name": name,
                    "type": ent_type,
                    "name_norm": name_norm,
                }
                name_index.setdefault(name_norm, key)
        else:
            key = str(uuid.uuid4())
            nodes[key] = {
                "id": str(uuid.uuid4()),
                "name": name,
                "type": ent_type,
                "name_norm": name_norm,
            }
            name_index.setdefault(name_norm, key)
    cleaned_relations = []
    for rel in relations:
        source = rel.get("source", "").strip()
        target = rel.get("target", "").strip()
        rel_type = rel.get("relation", "").strip() or "RELATED"
        if not source or not target:
            continue
        source_norm = _normalize_entity_name(source)
        target_norm = _normalize_entity_name(target)
        source_key = name_index.get(source_norm)
        if source_key is None:
            source_key = f"UNKNOWN::{source_norm}" if merge_entities else str(uuid.uuid4())
            nodes[source_key] = {
                "id": _entity_id(source, "UNKNOWN") if merge_entities else str(uuid.uuid4()),
                "name": source,
                "type": "UNKNOWN",
                "name_norm": source_norm,
            }
            name_index.setdefault(source_norm, source_key)
        target_key = name_index.get(target_norm)
        if target_key is None:
            target_key = f"UNKNOWN::{target_norm}" if merge_entities else str(uuid.uuid4())
            nodes[target_key] = {
                "id": _entity_id(target, "UNKNOWN") if merge_entities else str(uuid.uuid4()),
                "name": target,
                "type": "UNKNOWN",
                "name_norm": target_norm,
            }
            name_index.setdefault(target_norm, target_key)
        cleaned_relations.append(
            {
                "source_id": nodes[source_key]["id"],
                "target_id": nodes[target_key]["id"],
                "type": rel_type,
            }
        )
    return nodes, cleaned_relations


def ingest_to_neo4j(
    driver,
    nodes: Dict[str, Dict[str, str]],
    relations: List[Dict[str, str]],
    source_id: str,
    paragraph_id: int,
    paragraph_text: str | None = None,
    is_short: bool = False,
) -> None:
    with driver.session() as session:
        if paragraph_text:
            session.run(
                """
                MERGE (d:Document {id: $doc_id})
                SET d.name = $doc_name
                MERGE (p:Paragraph {source_id: $source_id, paragraph_id: $paragraph_id})
                SET p.text = $text,
                    p.short = $is_short
                MERGE (d)-[:HAS_PARAGRAPH]->(p)
                """,
                doc_id=source_id,
                doc_name=source_id,
                source_id=source_id,
                paragraph_id=paragraph_id,
                text=paragraph_text,
                is_short=is_short,
            )
        for node in nodes.values():
            session.run(
                """
                MERGE (e:Entity {id: $id})
                SET e.name = coalesce(e.name, $name),
                    e.type = $type,
                    e.name_norm = $name_norm
                """,
                id=node["id"],
                name=node["name"],
                type=node["type"],
                name_norm=node["name_norm"],
            )
            if paragraph_text:
                session.run(
                    """
                    MATCH (p:Paragraph {source_id: $source_id, paragraph_id: $paragraph_id})
                    MATCH (e:Entity {id: $id})
                    MERGE (p)-[r:HAS_ENTITY]->(e)
                    SET r.source_id = $source_id,
                        r.paragraph_id = $paragraph_id
                    """,
                    source_id=source_id,
                    paragraph_id=paragraph_id,
                    id=node["id"],
                )
        for rel in relations:
            session.run(
                """
                MATCH (s:Entity {id: $source_id})
                MATCH (t:Entity {id: $target_id})
                MERGE (s)-[r:RELATED {type: $type, source_id: $source_doc, paragraph_id: $paragraph_id}]->(t)
                """,
                source_id=rel["source_id"],
                target_id=rel["target_id"],
                type=rel["type"],
                source_doc=source_id,
                paragraph_id=paragraph_id,
            )


def ingest_to_neo4j_batch(driver, items: List[Dict[str, Any]]) -> None:
    paragraphs = []
    entities = []
    relations = []
    for item in items:
        paragraph_text = item.get("paragraph_text")
        if paragraph_text is not None:
            paragraphs.append(
                {
                    "doc_id": item["source_id"],
                    "doc_name": item["source_id"],
                    "source_id": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                    "text": paragraph_text,
                    "is_short": item.get("is_short", False),
                }
            )
        for node in item.get("nodes", {}).values():
            entities.append(
                {
                    "id": node["id"],
                    "name": node["name"],
                    "type": node["type"],
                    "name_norm": node["name_norm"],
                    "source_id": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                }
            )
        for rel in item.get("relations", []):
            relations.append(
                {
                    "source_id": rel["source_id"],
                    "target_id": rel["target_id"],
                    "type": rel["type"],
                    "source_doc": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                }
            )
    with driver.session() as session:
        if paragraphs:
            session.run(
                """
                UNWIND $paragraphs AS row
                MERGE (d:Document {id: row.doc_id})
                SET d.name = row.doc_name
                MERGE (p:Paragraph {source_id: row.source_id, paragraph_id: row.paragraph_id})
                SET p.text = row.text,
                    p.short = row.is_short
                MERGE (d)-[:HAS_PARAGRAPH]->(p)
                """,
                paragraphs=paragraphs,
            )
        if entities:
            session.run(
                """
                UNWIND $entities AS row
                MERGE (e:Entity {id: row.id})
                SET e.name = coalesce(e.name, row.name),
                    e.type = row.type,
                    e.name_norm = row.name_norm
                """,
                entities=entities,
            )
            session.run(
                """
                UNWIND $entities AS row
                MATCH (p:Paragraph {source_id: row.source_id, paragraph_id: row.paragraph_id})
                MATCH (e:Entity {id: row.id})
                MERGE (p)-[r:HAS_ENTITY]->(e)
                SET r.source_id = row.source_id,
                    r.paragraph_id = row.paragraph_id
                """,
                entities=entities,
            )
        if relations:
            session.run(
                """
                UNWIND $relations AS row
                MATCH (s:Entity {id: row.source_id})
                MATCH (t:Entity {id: row.target_id})
                MERGE (s)-[r:RELATED {type: row.type, source_id: row.source_doc, paragraph_id: row.paragraph_id}]->(t)
                """,
                relations=relations,
            )


def create_collection(client: QdrantClient, name: str, vector_size: int) -> None:
    try:
        client.get_collection(name)
        return
    except Exception:
        pass
    client.create_collection(
        collection_name=name,
        vectors_config=qmodels.VectorParams(size=vector_size, distance=qmodels.Distance.COSINE),
    )


def ingest_to_qdrant(
    client: QdrantClient,
    collection: str,
    paragraph: str,
    paragraph_id: int,
    embedder: SentenceTransformer,
    nodes: Dict[str, Dict[str, str]],
    source_id: str,
) -> None:
    node_lookup = {node["name"].lower(): node["id"] for node in nodes.values()}
    vector = embedder.encode([paragraph])[0]
    entity_ids = []
    lower_para = paragraph.lower()
    for name_lower, node_id in node_lookup.items():
        if name_lower and name_lower in lower_para:
            entity_ids.append(node_id)
    point = qmodels.PointStruct(
        id=str(uuid.uuid4()),
        vector=vector.tolist(),
        payload={
            "paragraph_id": paragraph_id,
            "source_id": source_id,
            "text": paragraph,
            "entity_ids": entity_ids,
        },
    )
    client.upsert(collection_name=collection, points=[point])


def _safe_source_id(base: Path, file_path: Path) -> str:
    rel = file_path.relative_to(base).as_posix()
    return rel.replace("/", "__").replace("\\", "__")


def _iter_input_files(folder: Path) -> List[Path]:
    files = []
    for path in folder.rglob("*"):
        if path.is_file() and path.suffix.lower() in {".pdf", ".txt", ".md", ".docx", ".pptx"}:
            files.append(path)
    return sorted(files)


def _read_input_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf_text(file_path)
    if suffix == ".docx":
        return read_docx_text(file_path)
    if suffix == ".pptx":
        return read_pptx_text(file_path)
    return read_text_file(file_path)


def process_text(
    raw_text: str,
    source_id: str,
    args: argparse.Namespace,
    driver,
    qdrant: QdrantClient,
    embedder: SentenceTransformer,
    nlp=None,
    gliner_model=None,
) -> None:
    if not raw_text:
        print(f"Skip empty input: {source_id}")
        return

    ingested_at = datetime.now(timezone.utc).isoformat()
    print(f"Source: {source_id} at {ingested_at}")
    print(f"Entity provider: {args.entity_provider}")

    if args.entity_provider == "spacy" and nlp is None:
        nlp = build_spacy_pipeline(args.spacy_model, ruler_json=args.ruler_json)
    if args.entity_provider == "gliner" and gliner_model is None:
        gliner_model = build_gliner_model(args.gliner_model_resolved)
    gliner_labels = parse_gliner_labels(args.gliner_labels)
    use_gliner_batch = args.entity_provider == "gliner" and not args.no_batch
    merge_entities = not args.no_entity_merge
    if args.no_batch:
        args.gliner_batch_size = 1
        args.neo4j_batch_size = 1

    paragraphs = split_paragraphs(raw_text, args.max_paragraph_chars)
    print(f"Chunked into {len(paragraphs)} paragraphs.")
    neo4j_batch: List[Dict[str, Any]] = []

    def flush_neo4j_batch() -> None:
        if not neo4j_batch:
            return
        print(f"Ingesting {len(neo4j_batch)} paragraphs to Neo4j (batch)...")
        ingest_to_neo4j_batch(driver, neo4j_batch)
        neo4j_batch.clear()
        print("Neo4j batch ingestion complete.")

    long_paragraphs: List[Tuple[int, str]] = []
    for idx, paragraph in enumerate(paragraphs):
        if not paragraph:
            continue
        if len(paragraph) < args.min_paragraph_chars:
            if args.skip_llm_short:
                print(
                    f"Paragraph {idx + 1}/{len(paragraphs)}: LLM skipped (len={len(paragraph)})"
                )
                neo4j_batch.append(
                    {
                        "source_id": source_id,
                        "paragraph_id": idx,
                        "paragraph_text": paragraph,
                        "is_short": True,
                        "nodes": {},
                        "relations": [],
                    }
                )
                if len(neo4j_batch) >= args.neo4j_batch_size:
                    flush_neo4j_batch()
                print("Ingesting to Qdrant...")
                ingest_to_qdrant(qdrant, args.collection, paragraph, idx, embedder, {}, source_id)
                print("Qdrant ingestion complete.")
            else:
                print(f"Paragraph {idx + 1}/{len(paragraphs)}: skipped (len={len(paragraph)})")
            continue
        long_paragraphs.append((idx, paragraph))

    if use_gliner_batch and long_paragraphs:
        batch_size = max(1, args.gliner_batch_size)
        for start in range(0, len(long_paragraphs), batch_size):
            batch_items = long_paragraphs[start : start + batch_size]
            batch_texts = [item[1] for item in batch_items]
            print(
                f"Extracting entities with GLiNER (batch {start + 1}-"
                f"{min(start + batch_size, len(long_paragraphs))})..."
            )
            batch_entities = extract_entities_gliner_batch(
                batch_texts,
                labels=gliner_labels,
                threshold=args.gliner_threshold,
                gliner_model=gliner_model,
            )
            for (idx, paragraph), entities in zip(batch_items, batch_entities, strict=True):
        nodes, relations = build_graph_components_from_entities(
            entities, [], merge_entities=merge_entities
        )
                print(f"Paragraph {idx + 1}/{len(paragraphs)}: extracted {len(nodes)} entities.")
                neo4j_batch.append(
                    {
                        "source_id": source_id,
                        "paragraph_id": idx,
                        "paragraph_text": paragraph,
                        "is_short": False,
                        "nodes": nodes,
                        "relations": relations,
                    }
                )
                if len(neo4j_batch) >= args.neo4j_batch_size:
                    flush_neo4j_batch()
                print("Ingesting to Qdrant...")
                ingest_to_qdrant(qdrant, args.collection, paragraph, idx, embedder, nodes, source_id)
                print("Qdrant ingestion complete.")
    else:
        for idx, paragraph in long_paragraphs:
            print(f"Paragraph {idx + 1}/{len(paragraphs)}: extracting entities/relations...")
            nodes, relations = build_graph_components(
                paragraph,
                args.entity_provider,
                nlp=nlp,
                gliner_model=gliner_model,
                gliner_labels=gliner_labels,
                gliner_threshold=args.gliner_threshold,
                merge_entities=merge_entities,
            )
            print(f"Extracted {len(nodes)} entities, {len(relations)} relations.")
            neo4j_batch.append(
                {
                    "source_id": source_id,
                    "paragraph_id": idx,
                    "paragraph_text": paragraph,
                    "is_short": False,
                    "nodes": nodes,
                    "relations": relations,
                }
            )
            if len(neo4j_batch) >= args.neo4j_batch_size:
                flush_neo4j_batch()
            print("Ingesting to Qdrant...")
            ingest_to_qdrant(qdrant, args.collection, paragraph, idx, embedder, nodes, source_id)
            print("Qdrant ingestion complete.")
    flush_neo4j_batch()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pdf", help="Path to PDF")
    parser.add_argument("--text-file", help="Path to UTF-8 text file")
    parser.add_argument("--raw-text", help="Raw text input")
    parser.add_argument("--folder", help="Folder to scan for .pdf/.txt files (recursive)")
    parser.add_argument("--source-id", default=None, help="Custom source identifier")
    parser.add_argument("--collection", default="graphrag_entities")
    parser.add_argument("--embedding-model", default=None)
    parser.add_argument("--max-paragraph-chars", type=int, default=1200)
    parser.add_argument("--min-paragraph-chars", type=int, default=150)
    parser.add_argument(
        "--skip-llm-short",
        action="store_true",
        help="Skip LLM extraction for short paragraphs but still store in Qdrant",
    )
    parser.add_argument(
        "--entity-provider",
        choices=["spacy", "gliner", "gemini", "langextract"],
        default="gliner",
        help="Entity extraction provider",
    )
    parser.add_argument("--spacy-model", default="en_core_web_sm")
    parser.add_argument("--ruler-json", default=None, help="Path to EntityRuler JSON patterns")
    parser.add_argument(
        "--gliner-model-name",
        default=os.getenv("GLINER_MODEL_NAME", "urchade/gliner_mediumv2"),
        help="GLiNER model name (HuggingFace).",
    )
    parser.add_argument(
        "--gliner-model-path",
        default=os.getenv("GLINER_MODEL_PATH"),
        help="Local path to a GLiNER model directory.",
    )
    parser.add_argument(
        "--gliner-model",
        default=None,
        help="Deprecated: use --gliner-model-name or --gliner-model-path.",
    )
    parser.add_argument(
        "--gliner-labels",
        default="PERSON,ORG,PRODUCT,GPE,DATE,TECH,CRYPTO,STANDARD",
        help="Comma-separated labels or path to a text file for GLiNER",
    )
    parser.add_argument("--gliner-threshold", type=float, default=0.3)
    parser.add_argument(
        "--gliner-batch-size",
        type=int,
        default=8,
        help="Batch size for GLiNER entity extraction.",
    )
    parser.add_argument(
        "--no-entity-merge",
        action="store_true",
        help="Do not merge entities across paragraphs (use per-paragraph IDs).",
    )
    parser.add_argument(
        "--no-batch",
        action="store_true",
        help="Disable batching for GLiNER and Neo4j (sequential processing).",
    )
    parser.add_argument("--langextract-model-id", default=os.getenv("LANGEXTRACT_MODEL_ID"))
    parser.add_argument("--langextract-model-url", default=os.getenv("LANGEXTRACT_MODEL_URL"))
    parser.add_argument("--llm-debug", action="store_true", help="Print raw LLM output")
    parser.add_argument("--llm-retry-count", type=int, default=None)
    parser.add_argument("--llm-retry-backoff", type=float, default=None)
    parser.add_argument(
        "--neo4j-batch-size",
        type=int,
        default=50,
        help="Number of paragraphs to write per Neo4j batch.",
    )
    parser.add_argument("--neo4j-uri", default=os.getenv("NEO4J_URI", "bolt://localhost:7687"))
    parser.add_argument("--neo4j-user", default=os.getenv("NEO4J_USER", "neo4j"))
    parser.add_argument("--neo4j-pass", default=os.getenv("NEO4J_PASS", "password"))
    parser.add_argument("--qdrant-url", default=os.getenv("QDRANT_URL"))
    parser.add_argument("--qdrant-host", default=os.getenv("QDRANT_HOST", "localhost"))
    parser.add_argument("--qdrant-port", type=int, default=int(os.getenv("QDRANT_PORT", "6333")))
    parser.add_argument("--qdrant-api-key", default=os.getenv("QDRANT_KEY"))
    args = parser.parse_args()

    _load_env()

    if args.gliner_model_path:
        gliner_model_choice = args.gliner_model_path
    elif args.gliner_model:
        gliner_model_choice = args.gliner_model
    else:
        gliner_model_choice = args.gliner_model_name
    args.gliner_model_resolved = gliner_model_choice

    inputs = [bool(args.pdf), bool(args.text_file), bool(args.raw_text), bool(args.folder)]
    if sum(inputs) != 1:
        raise SystemExit("Provide exactly one of --pdf, --text-file, --raw-text, or --folder.")

    if args.llm_debug:
        os.environ["LLM_DEBUG"] = "1"
        os.environ.setdefault("LLM_DEBUG_LOG_PATH", "logs/langextract_raw.log")
    if args.llm_retry_count is not None:
        os.environ["LLM_RETRY_COUNT"] = str(args.llm_retry_count)
    if args.llm_retry_backoff is not None:
        os.environ["LLM_RETRY_BACKOFF_SECONDS"] = str(args.llm_retry_backoff)
    _set_langextract_overrides(args.langextract_model_id, args.langextract_model_url)
    model_name, local_files_only = resolve_embedding_model(
        args.embedding_model, "sentence-transformers/all-MiniLM-L6-v2"
    )
    embedder = SentenceTransformer(model_name, local_files_only=local_files_only)

    if args.qdrant_url:
        qdrant = QdrantClient(url=args.qdrant_url, api_key=args.qdrant_api_key)
    else:
        qdrant = QdrantClient(
            host=args.qdrant_host, port=args.qdrant_port, api_key=args.qdrant_api_key
        )
    create_collection(qdrant, args.collection, vector_size=embedder.get_sentence_embedding_dimension())

    driver = GraphDatabase.driver(args.neo4j_uri, auth=(args.neo4j_user, args.neo4j_pass))

    if args.folder:
        folder_path = Path(args.folder)
        if not folder_path.exists():
            raise FileNotFoundError(folder_path)
        file_paths = _iter_input_files(folder_path)
        if not file_paths:
            raise SystemExit("No supported files found in folder.")
        print(f"Found {len(file_paths)} files in folder.")
        shared_nlp = None
        shared_gliner = None
        if args.entity_provider == "spacy":
            shared_nlp = build_spacy_pipeline(args.spacy_model, ruler_json=args.ruler_json)
        if args.entity_provider == "gliner":
            shared_gliner = build_gliner_model(args.gliner_model_resolved)
        for idx, file_path in enumerate(file_paths, start=1):
            print(f"[{idx}/{len(file_paths)}] Processing: {file_path}")
            if args.source_id:
                source_id = f"{args.source_id}__{_safe_source_id(folder_path, file_path)}"
            else:
                source_id = _safe_source_id(folder_path, file_path)
            raw_text = _read_input_text(file_path)
            process_text(
                raw_text,
                source_id,
                args,
                driver,
                qdrant,
                embedder,
                nlp=shared_nlp,
                gliner_model=shared_gliner,
            )
        return

    if args.pdf:
        pdf_path = Path(args.pdf)
        if not pdf_path.exists():
            raise FileNotFoundError(pdf_path)
        raw_text = read_pdf_text(pdf_path)
        source_id = args.source_id or pdf_path.stem
        process_text(raw_text, source_id, args, driver, qdrant, embedder)
        return

    if args.text_file:
        text_path = Path(args.text_file)
        if not text_path.exists():
            raise FileNotFoundError(text_path)
        raw_text = read_text_file(text_path)
        source_id = args.source_id or text_path.stem
        process_text(raw_text, source_id, args, driver, qdrant, embedder)
        return

    raw_text = args.raw_text.strip()
    source_id = args.source_id or "raw_text"
    process_text(raw_text, source_id, args, driver, qdrant, embedder)


if __name__ == "__main__":
    main()
